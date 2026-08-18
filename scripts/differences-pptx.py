"""Convert DIFFERENCES.md into differences.pptx, one slide per building block group.

For every differing field a table "PFS | Current <field> | Proposed <field>" is
created. The Proposed column is derived from the group's manual Proposal section:
explicit texts and quoted rewordings are shown directly, "X match Y" directives
are resolved to Y's current text, removals become a DELETE cell, "No change"
repeats the current text. Proposed changes are shown as a diff against the
row's current text (red strikethrough = removed, green = added); plain =
unchanged, red DELETE = deletion, gray italic = directive that could not be
interpreted (shown verbatim). Raw Proposal and Internal notes become speaker notes.

Usage: python scripts/differences-pptx.py
"""

import importlib.util
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
spec = importlib.util.spec_from_file_location("cbb", HERE / "compare-building-blocks.py")
cbb = importlib.util.module_from_spec(spec)
spec.loader.exec_module(cbb)

OUT = cbb.ROOT / "differences.pptx"

GREEN = RGBColor(0x2E, 0x7D, 0x32)
GRAY = RGBColor(0x80, 0x80, 0x80)
RED = RGBColor(0xB0, 0x00, 0x00)
BLUE = RGBColor(0x15, 0x65, 0xC0)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.5)
COL_W = [Inches(1.0), Inches(5.66), Inches(5.66)]
PFS_RE = r"\b(?:NLSR|AR|SR|ST|CB|GSLC|INSAR|NRB|ORB|POL)\b"


def clean(text):
    if text is None:
        return "(not set)"
    if isinstance(text, bool):
        return str(text).lower()
    if isinstance(text, list):
        return "\n\n".join(text) if text else "(no notes)"
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text)


def field_title(key):
    return cbb.field_label(key).replace(" — ", " ").title()


# --------------------------------------------------------------------------
# Proposal interpretation
# --------------------------------------------------------------------------

def parse_proposal(text):
    """Split a Proposal section into entries: (field key, scope, fenced text, directive)."""
    lines = [l for l in text.split("\n") if not l.startswith("> ")]
    entries = []
    i = 0
    while i < len(lines):
        m = re.match(r"\*\*(.+?)\*\*(?:\s*\([^)]*\))?\s*:?\s*(.*)$", lines[i])
        if m:
            label, directive = m.group(1), m.group(2).strip() or None
            scope = None
            sm = re.match(r"(.+?)\s*\(([^)]*)\)\s*$", label)
            if sm:
                label = sm.group(1)
                scope = set(re.findall(PFS_RE, sm.group(2)))
            fenced = None
            j = i + 1
            while j < len(lines) and not lines[j].strip():
                j += 1
            if j < len(lines) and re.fullmatch(r"`{3,}", lines[j]):
                fence, start = lines[j], j + 1
                j += 1
                while lines[j] != fence:
                    j += 1
                fenced = "\n".join(lines[start:j])
                i = j
            entries.append({"key": cbb.label_to_key(label), "scope": scope,
                            "text": fenced, "directive": directive})
        i += 1
    return entries


def source_value(code, key, files, fields, usage):
    for f in files:
        if code in usage.get(f, []):
            return fields[f][key]
    return None


def resolve_directive(directive, key, row_pfs, row_value, files, fields, usage):
    """Interpret one directive for one table row; returns (value, style) or None."""
    sentences = re.split(r"(?<=[.\"])\s+(?=[A-Z*])", directive)
    for s in sentences:
        low = s.lower()
        mm = re.search(rf"(?:^|\b)(.*?)match(?:es)? ({PFS_RE})"
                       rf"( wording at threshold| threshold)?\b(.*)$", s, re.I)
        codes_before = set(re.findall(PFS_RE, mm.group(1))) if mm else set()
        if mm:
            scope = codes_before or None
            src = mm.group(2)
            if scope and not (scope & row_pfs):
                if src in row_pfs:
                    return clean(row_value), "same"
                continue
            src_key = "threshold.description" if mm.group(3) else key
            tail = mm.group(4) or ""
            quoted = re.search(r'"([^"]+)"', tail)
            if quoted and ("phrasing" in low or tail.strip().startswith('"')):
                value = quoted.group(1)
            else:
                value = source_value(src, src_key, files, fields, usage)
                if key.endswith(".notes") and (value is None or value == []):
                    return "DELETE NOTE", "delete"
                value = clean(value)
                extra = re.search(r"\b(BUT .+)$", tail)
                if extra:
                    value += f"\n[{extra.group(1)}]"
            return value, ("same" if value == clean(row_value) else "changed")
        codes = set(re.findall(PFS_RE, s))
        applies = not codes or codes & row_pfs
        if not applies:
            continue
        if re.search(r"\b(remove|delete)\b", low):
            return ("DELETE NOTE" if key.endswith(".notes") else "DELETE"), "delete"
        if "no change" in low or "keep" in low:
            return clean(row_value), "same"
        quoted = re.search(r'"([^"]+)"', s)
        if quoted:
            value = quoted.group(1)
            return value, ("same" if value == clean(row_value) else "changed")
        if low.rstrip(".") == "as threshold":
            return "As threshold", ("same" if clean(row_value) == "As threshold" else "changed")
        return s, "raw"
    return None


def proposed_cell(entries, key, row_pfs, row_value, files, fields, usage):
    candidates = [e for e in entries if e["key"] == key
                  and (e["scope"] is None or e["scope"] & row_pfs)]
    candidates.sort(key=lambda e: e["scope"] is None)  # scoped entries first
    for e in candidates:
        if e["text"] is not None:
            value = e["text"]
            return clean(value), ("same" if clean(value) == clean(row_value) else "changed")
        result = resolve_directive(e["directive"], key, row_pfs, row_value,
                                   files, fields, usage)
        if result:
            return result
    return None


# --------------------------------------------------------------------------
# Diff coloring of the Current column, mirroring the markdown highlighting
# --------------------------------------------------------------------------

def diff_spans(base_text, cur_text):
    """Spans (text, style) for a variant compared to the baseline text."""
    if base_text == cur_text:
        return [(cur_text, None)]
    spans = []
    if cbb.similarity(base_text, cur_text) >= cbb.DIFF_THRESHOLD:
        for seg in cbb.diff_segments(base_text, cur_text):
            if seg[0] == "equal":
                spans.append((seg[1], None))
            else:
                if seg[1]:
                    spans.append((seg[1], "del"))
                if seg[2]:
                    spans.append((seg[2], "ins"))
    else:  # too different for a diff: highlight the shared passages instead
        for seg in cbb.diff_segments(base_text, cur_text):
            if seg[0] == "equal":
                spans.append((seg[1], "sim" if len(seg[1]) >= cbb.MIN_EQUAL else None))
            else:
                spans.append((seg[2], None))
    return [s for s in spans if s[0]]


# --------------------------------------------------------------------------
# Slide building
# --------------------------------------------------------------------------

def cell_len(spans):
    return sum(len(t) for t, s in spans if s != "del")


def est_height(rows):
    h = 0.32  # header
    for cells in rows:
        lines = max(sum(len(part) // 78 + 1
                        for part in "".join(t for t, _ in c).split("\n")) for c in cells)
        h += 0.11 + lines * 0.165
    return h


def add_title(slide, text):
    box = slide.shapes.add_textbox(MARGIN, Inches(0.25), SLIDE_W - 2 * MARGIN, Inches(0.7))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    return box


def write_cell(cell, spans, bold_default=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame
    para = tf.paragraphs[0]
    for text, style in spans:
        for i, part in enumerate((text or "").split("\n")):
            if i > 0:
                para = tf.add_paragraph()
            if not part:
                continue
            run = para.add_run()
            run.text = part
            font = run.font
            font.size = Pt(10)
            if style == "del":
                font.color.rgb = RED
                font._rPr.set("strike", "sngStrike")
            elif style in ("ins", "changed"):
                font.color.rgb = GREEN
                font.bold = True
            elif style == "sim":
                font.color.rgb = BLUE
            elif style == "delete":
                font.color.rgb = RED
                font.bold = True
            elif style in ("raw", "note"):
                font.color.rgb = GRAY
                font.italic = True
            elif bold_default:
                font.bold = True


def add_table(slide, y, headers, rows):
    shape = slide.shapes.add_table(len(rows) + 1, 3, MARGIN, y,
                                   SLIDE_W - 2 * MARGIN, Inches(0.4))
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    for c, w in enumerate(COL_W):
        table.columns[c].width = w
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        cell.text = text
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
    for r, cells in enumerate(rows, start=1):
        for c, spans in enumerate(cells):
            write_cell(table.cell(r, c), spans, bold_default=(c == 0))
    return shape


def joint_value(files, fields, usage, entries, key):
    """Proposed end state of one field across the group; True if it still differs."""
    vals = []
    for f in files:
        cur = fields[f].get(key)
        prop = proposed_cell(entries, key, set(usage.get(f, [])), cur, files, fields, usage)
        if prop is None or prop[1] == "raw":
            v = cur if isinstance(cur, bool) else clean(cur)
        elif prop[1] == "delete":
            v = ""
        else:
            v = prop[0]
        if v in ("(not set)", "(no notes)"):
            v = ""
        vals.append(v)
    counts = {}
    for v in vals:
        counts[v] = counts.get(v, 0) + 1
    value = max(counts, key=counts.get)
    return value, len(counts) > 1


def joint_yaml_lines(files, fields, usage, entries, levels):
    """The proposed joint building block as plain YAML lines."""
    lines = []

    def scalar(indent, key, value, todo):
        suffix = "  # TODO: differs per PFS" if todo else ""
        if "\n" in value:
            lines.append(f"{indent}{key}: |-{suffix}")
            for l in value.split("\n"):
                lines.append(f"{indent}  {l}")
        else:
            lines.append(f"{indent}{key}: {value}{suffix}")

    for key in ("title", "description"):
        value, todo = joint_value(files, fields, usage, entries, key)
        if value:
            scalar("", key, value, todo)
    lines.append("requirements:")
    for lvl in cbb.LEVEL_ORDER:
        if lvl not in levels:
            continue
        desc, d_todo = joint_value(files, fields, usage, entries, f"{lvl}.description")
        notes, n_todo = joint_value(files, fields, usage, entries, f"{lvl}.notes")
        opt, o_todo = joint_value(files, fields, usage, entries, f"{lvl}.optional")
        if not (desc or notes or opt):
            continue
        lines.append(f"  {lvl}:")
        if desc:
            scalar("    ", "description", desc, d_todo)
        if notes:
            lines.append("    notes:" + ("  # TODO: differs per PFS" if n_todo else ""))
            for note in notes.split("\n\n"):
                if "\n" in note:
                    lines.append("      - |-")
                    for l in note.split("\n"):
                        lines.append(f"        {l}")
                else:
                    lines.append(f"      - {note}")
        if opt is True:
            lines.append("    optional: true" + ("  # TODO: differs per PFS" if o_todo else ""))
    return lines


def add_proposal_slides(prs, blank, title, lines):
    """Slides with the proposed joint building block as editable plain YAML."""
    per_slide = []
    current, height = [], 0
    for line in lines:
        h = max(1, len(line) // 115 + 1)
        if current and height + h > 28:
            per_slide.append(current)
            current, height = [], 0
        current.append(line)
        height += h
    per_slide.append(current)
    for n, slide_lines in enumerate(per_slide):
        slide = prs.slides.add_slide(blank)
        add_title(slide, title + " — Proposal" + (" (cont.)" if n else ""))
        box = slide.shapes.add_textbox(MARGIN, Inches(1.1), SLIDE_W - 2 * MARGIN, Inches(6.2))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for line in slide_lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = line
            run.font.size = Pt(11)
            run.font.name = "Consolas"


def main():
    groups = cbb.collect_groups()
    usage, overrides, cat_overrides, types = cbb.collect_pfs()
    groups, overrides = cbb.drop_sar_only(groups, overrides, usage, types)
    manual, _states = cbb.extract_manual(cbb.OUTPUT.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    # title slide
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(MARGIN, Inches(2.6), SLIDE_W - 2 * MARGIN, Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Differences between similar building blocks"
    run.font.size = Pt(40)
    run.font.bold = True
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("CEOS-ARD Product Family Specifications — current wording per PFS "
                "and alignment proposals from the review feedback")
    run.font.size = Pt(18)
    run.font.color.rgb = GRAY

    # scope slide: context from the DIFFERENCES.md introduction
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Scope")
    box = slide.shapes.add_textbox(MARGIN, Inches(1.2), SLIDE_W - 2 * MARGIN, Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    scope = [
        "Building blocks that exist in several variants (e.g. time-ar/-sar/-sr/-st) are "
        "compared; only the differences are shown — fields that are identical across all "
        "variants of a group are omitted.",
        "The groups are split by the type of the PFS that use them: used by both SAR and "
        "Optical PFS, or used only by Optical PFS.",
        "Building blocks used only by SAR PFS are omitted (as are PFS overrides that only "
        "affect them): their differences are intentional, as the SAR PFS were a combined "
        "document before being split.",
        "Compared fields: title, description, and per requirement level (threshold/goal/"
        "image) the description, the optional flag, and the notes.",
        "Ignored fields: id, dependencies, glossary, references, changes, history, remarks.",
        "Not compared: the measurement/measurand and backscatter building blocks "
        "(intentionally different per product).",
    ]
    first = True
    for bullet in scope:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + bullet
        run.font.size = Pt(16)

    # legend slide
    slide = prs.slides.add_slide(blank)
    add_title(slide, "How to read the tables")
    box = slide.shapes.add_textbox(MARGIN, Inches(1.2), SLIDE_W - 2 * MARGIN, Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    legend = [
        [("One slide per building block group. For every field that differs between the "
          "variants, a table shows the current text per PFS and the proposed alignment.", None)],
        [("", None)],
        [("Current column — differences to the row marked ", None), ("(baseline)", "note"),
         (" are highlighted:", None)],
        [("    ", None), ("red struck-through", "del"),
         (" text only appears in the baseline variant", None)],
        [("    ", None), ("green", "ins"), (" text only appears in this variant", None)],
        [("    ", None), ("blue passages", "sim"),
         (" are shared with the baseline — used where texts are too different for a diff", None)],
        [("    ", None), ("(optional) / (required)", "note"),
         (" shows the optional flag where it differs between the variants", None)],
        [("    (not set) = the field or the whole requirement level is absent in that file", None)],
        [("", None)],
        [("Proposed column — proposed changes relative to the row's current text:", None)],
        [("    ", None), ("red struck-through", "del"), (" text is removed, ", None),
         ("green", "ins"), (" text is added by the proposal; plain text = unchanged", None)],
        [("    ", None), ("DELETE NOTE", "delete"), (" = proposed removal", None)],
        [("    ", None), ("gray italic", "raw"),
         (" = proposal that could not be interpreted automatically, shown verbatim", None)],
        [("", None)],
        [("⚠ in the slide title: the files changed after the feedback was written — "
          "review whether the proposal still applies.", None)],
        [("Internal notes are shown below the tables; the raw proposal text is in the "
          "slide's speaker notes.", None)],
        [("Each group is followed by a Proposal slide showing the proposed building blocks "
          "in YAML form (only the fields that change; proposed values in ", None),
         ("green", "ins"), (") for editing and capturing the agreement during the review.", None)],
    ]
    first = True
    for spans in legend:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        for text, style in spans:
            run = p.add_run()
            run.text = text
            font = run.font
            font.size = Pt(14)
            if style == "del":
                font.color.rgb = RED
                font._rPr.set("strike", "sngStrike")
            elif style in ("ins", "changed"):
                font.color.rgb = GREEN
                font.bold = True
            elif style == "sim":
                font.color.rgb = BLUE
            elif style == "delete":
                font.color.rgb = RED
                font.bold = True
            elif style in ("raw", "note"):
                font.color.rgb = GRAY
                font.italic = True

    for label, files in groups.items():
        docs = {f: cbb.load(f) for f in files}
        levels = cbb.group_levels(docs)
        fields = {f: cbb.extract_fields(doc, levels) for f, doc in docs.items()}
        titles = {fields[f]["title"] for f in files}
        heading = titles.pop() if len(titles) == 1 else f"`{label}`"
        entries_raw = manual.get(heading, {})
        proposal_text = entries_raw.get("Proposal", "")
        internal_notes = entries_raw.get("Internal notes", "")
        entries = parse_proposal(proposal_text) if proposal_text else []
        stale = cbb.STALE_WARNING in proposal_text or cbb.STALE_WARNING in internal_notes
        title = heading.replace("`", "") + (" ⚠" if stale else "")

        diff_keys = [k for k in cbb.field_keys(levels)
                     if not all(fields[f][k] == fields[files[0]][k] for f in files)]
        # the optional flag is shown as annotation in the description table, not as own table
        opt_levels = {k.split(".")[0] for k in diff_keys if k.endswith(".optional")}
        tables = []
        extra_lines = []
        for level in sorted(opt_levels):
            if f"{level}.description" not in diff_keys:
                states_ = {}
                for f in files:
                    label_ = ", ".join(usage.get(f, [])) or f.rsplit("/", 1)[1]
                    states_.setdefault("optional" if fields[f][f"{level}.optional"] else "required",
                                       []).append(label_)
                extra_lines.append(f"{level.capitalize()}: "
                                   + "; ".join(f"{k} for {', '.join(v)}" for k, v in states_.items()))
        for key in diff_keys:
            if key.endswith(".optional"):
                continue
            variants = []  # [files, value] with identical values grouped
            for f in files:
                v = fields[f][key]
                for entry in variants:
                    if entry[1] == v:
                        entry[0].append(f)
                        break
                else:
                    variants.append([[f], v])
            diffable = [v for v in variants
                        if (isinstance(v[1], str) or isinstance(v[1], list)) and v[1]]
            base_variant = max(diffable, key=lambda v: len(v[0])) if diffable else None
            base_disp = clean(base_variant[1]) if base_variant else None
            level = key.split(".")[0] if "." in key else None
            rows = []
            for f in files:
                pfs_list = usage.get(f, [])
                row_pfs = set(pfs_list)
                pfs_spans = [(", ".join(pfs_list) or f.rsplit("/", 1)[1], None)]
                v = fields[f][key]
                if base_variant and f in base_variant[0]:
                    cur_spans = [(clean(v), None)]
                    if f == base_variant[0][0] and len(variants) > 1:
                        pfs_spans.append(("\n(baseline)", "note"))
                elif base_variant and (isinstance(v, (str, list)) and v):
                    cur_spans = diff_spans(base_disp, clean(v))
                else:
                    cur_spans = [(clean(v), None)]
                if level in opt_levels and key == f"{level}.description":
                    flag = fields[f][f"{level}.optional"]
                    cur_spans = cur_spans + [(f"\n({'optional' if flag else 'required'})", "note")]
                prop = proposed_cell(entries, key, row_pfs, fields[f][key],
                                     files, fields, usage)
                if prop is None:
                    prop_spans = [("", None)]
                elif prop[1] == "changed":
                    # highlight the proposed changes relative to this row's current text
                    if isinstance(v, (str, list)) and v:
                        prop_spans = diff_spans(clean(v), prop[0])
                    else:
                        prop_spans = [(prop[0], "ins")]
                else:
                    prop_spans = [prop]
                rows.append([pfs_spans, cur_spans, prop_spans])
            tables.append((field_title(key), rows))

        slide = prs.slides.add_slide(blank)
        add_title(slide, title)
        y = 1.1
        for tlabel, rows in tables:
            h = est_height(rows)
            if y + h > 7.2 and y > 1.1:
                slide = prs.slides.add_slide(blank)
                add_title(slide, title + " (cont.)")
                y = 1.1
            add_table(slide, Inches(y), ["PFS", f"Current {tlabel}", f"Proposed {tlabel}"], rows)
            y += h + 0.3
        for line in extra_lines:
            box = slide.shapes.add_textbox(MARGIN, Inches(y), SLIDE_W - 2 * MARGIN, Inches(0.3))
            box.text_frame.word_wrap = True
            box.text_frame.text = line
            for p in box.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.italic = True
            y += 0.35
        # internal notes, without the stale warning (already conveyed by the title icon)
        notes_clean = "\n".join(l for l in internal_notes.split("\n")
                                if not l.startswith("> ")).strip()
        if notes_clean:
            h = 0.2 + sum(len(l) // 110 + 1 for l in notes_clean.split("\n")) * 0.24
            if y + h > 7.3 and y > 1.1:
                slide = prs.slides.add_slide(blank)
                add_title(slide, title + " (cont.)")
                y = 1.1
            box = slide.shapes.add_textbox(MARGIN, Inches(y), SLIDE_W - 2 * MARGIN, Inches(h))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Internal notes: "
            run.font.size = Pt(12)
            run.font.bold = True
            first = True
            for line in notes_clean.split("\n"):
                if not first:
                    p = tf.add_paragraph()
                first = False
                run = p.add_run()
                run.text = line
                run.font.size = Pt(12)
                run.font.italic = True
            y += h
        notes = []
        if proposal_text:
            notes.append("PROPOSAL:\n" + proposal_text)
        if internal_notes:
            notes.append("INTERNAL NOTES:\n" + internal_notes)
        if notes:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(notes)
        if proposal_text:
            add_proposal_slides(prs, blank, title.replace(" ⚠", ""),
                                joint_yaml_lines(files, fields, usage, entries, levels))

    if cat_overrides:
        slide = prs.slides.add_slide(blank)
        add_title(slide, "PFS-specific overrides of requirement categories")
        rows = []
        for pfs, ref, mode, data in cat_overrides:
            for key, value in cbb.flatten_override(data)[0]:
                verb = "appends to" if mode == "append" else "replaces"
                rows.append([[(pfs, None)],
                             [(f"{ref} — {verb} {cbb.field_label(key).lower()}", None)],
                             [(clean(value), None)]])
        add_table(slide, Inches(1.1), ["PFS", "Category", "Text"], rows)
        box = slide.shapes.add_textbox(MARGIN, Inches(1.3 + est_height(rows)),
                                       SLIDE_W - 2 * MARGIN, Inches(1.0))
        box.text_frame.word_wrap = True
        box.text_frame.text = ('Proposal (from feedback): Append "DOI-landing page" verbiage to the '
                               'General Metadata and Per-Pixel Metadata sections for SR, ST and NLSR.')
        for p in box.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = GREEN

    prs.save(OUT)
    out_name = OUT.relative_to(cbb.ROOT) if OUT.is_relative_to(cbb.ROOT) else OUT.name
    print(f"Wrote {out_name} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
